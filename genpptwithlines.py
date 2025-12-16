from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# --- SETUP ---
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(10) # Increased height slightly for better spacing
slide = prs.slides.add_slide(prs.slide_layouts[6]) 

# Dictionary to store shape objects so we can connect them later
shape_map = {}

# --- HELPER FUNCTIONS ---

def create_block(left, top, width, height, color_rgb, title, content_list, key_name=None, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    """Creates a block and saves it to the shape_map"""
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    
    # Fill Color
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color_rgb
    shape.line.fill.background() # No outline for cleaner look

    # Text Styling
    text_frame = shape.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT

    for item in content_list:
        p = text_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(8)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.LEFT

    # Store shape for connecting later
    # Use title as key unless a specific key_name is provided
    store_key = key_name if key_name else title
    shape_map[store_key] = shape
    return shape

def add_connector(from_key, to_key, type=MSO_CONNECTOR.ELBOW):
    """Draws a line between two named shapes"""
    if from_key in shape_map and to_key in shape_map:
        source = shape_map[from_key]
        dest = shape_map[to_key]
        connector = slide.shapes.add_connector(type, 0, 0, 0, 0)
        
        # Connect the shapes
        # idx 0=top, 1=right, 2=bottom, 3=left (usually)
        connector.begin_connect(source, 2) # Start from bottom of source
        connector.end_connect(dest, 0)     # End at top of dest
        
        # Style the line
        line = connector.line
        line.color.rgb = RGBColor(89, 89, 89)
        line.width = Pt(1.5)

def add_side_connector(from_key, to_key):
    """Draws a line from Right of Source to Left of Dest"""
    if from_key in shape_map and to_key in shape_map:
        source = shape_map[from_key]
        dest = shape_map[to_key]
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, 0, 0, 0, 0)
        connector.begin_connect(source, 1) # Right
        connector.end_connect(dest, 3)     # Left
        connector.line.color.rgb = RGBColor(89, 89, 89)
        connector.line.width = Pt(1.5)

# --- COLORS ---
COL_BLUE_DARK = RGBColor(47, 85, 151)
COL_BLUE_LIGHT = RGBColor(68, 114, 196)
COL_PURPLE = RGBColor(112, 48, 160)
COL_TEAL = RGBColor(0, 176, 240)
COL_ORANGE = RGBColor(237, 125, 49)
COL_YELLOW = RGBColor(255, 192, 0)
COL_GREEN = RGBColor(146, 208, 80)
COL_RED = RGBColor(255, 100, 100)
COL_GREY_BG = RGBColor(242, 242, 242)
COL_DARK_GREY = RGBColor(89, 89, 89)

# --- HEADER ---
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(15), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "ICMP – Enterprise Content Management Capabilities Matrix"
p.font.bold = True
p.font.size = Pt(24)
p.font.color.rgb = COL_BLUE_DARK

# --- 1. TOP ROW (Ingestion) ---
y_top = 1.2
w_small = 1.5
h_small = 0.6
gap = 0.2
start_x = 0.5

sources = [
    ("ICMP UI", COL_BLUE_LIGHT),
    ("Apigee Gateway", COL_PURPLE),
    ("NAS Storage", COL_TEAL),
    ("Kafka Streams", COL_ORANGE),
    ("Email Ingestion", COL_YELLOW),
    ("Fax", COL_GREEN),
    ("Capture/KOFAX", COL_RED),
    ("Scanning", COL_RED)
]

for i, (name, col) in enumerate(sources):
    create_block(start_x + (i * (w_small + gap)), y_top, w_small, h_small, col, name, [], key_name=name)

# --- 2. FIREWALL ---
create_block(6.5, 2.2, 3, 0.4, COL_DARK_GREY, "PAA Firewalls", [], shape_type=MSO_SHAPE.RECTANGLE)

# --- 3. MAIN CONTAINER ---
bg_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.0), Inches(13), Inches(4.8))
bg_shape.fill.solid()
bg_shape.fill.fore_color.rgb = COL_GREY_BG
bg_shape.line.color.rgb = COL_DARK_GREY
bg_shape.text_frame.text = "ICMP Core Processing & Services Matrix"
bg_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
bg_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
bg_shape.text_frame.paragraphs[0].font.bold = True
# Map the background so we can connect to it roughly
shape_map["CORE_CONTAINER"] = bg_shape 

# --- 4. INSIDE BLOCKS ---

# Column 1
x1 = 0.7; w1 = 2.4
create_block(x1, 3.5, w1, 1.6, COL_BLUE_DARK, "API Layer", ["Ingestion", "Retrieval", "Search"], key_name="API")
create_block(x1, 5.3, w1, 1.1, COL_BLUE_LIGHT, "UI Module", ["Upload", "Re-index"], key_name="UI")
create_block(x1, 6.6, w1, 1.0, COL_BLUE_LIGHT, "IVaaS", ["Viewer", "Rendering"], key_name="IVaaS")

# Column 2
x2 = 3.3; w2 = 2.4
create_block(x2, 3.5, w2, 2.0, COL_BLUE_LIGHT, "ECM AI", ["Classification", "Extraction"], key_name="AI")
create_block(x2, 5.7, w2, 0.8, COL_PURPLE, "Transformation", ["Mapping", "Normalization"], key_name="TRANS")
create_block(x2, 6.7, w2, 0.9, COL_TEAL, "Conversion", ["PDF/TIFF", "Split/Merge"], key_name="CONV")

# Column 3
x3 = 5.9; w3 = 2.2
create_block(x3, 3.5, w3, 1.2, RGBColor(180, 50, 100), "Content Data", ["OCR/ICR", "Data lift"], key_name="DATA")
create_block(x3, 4.9, w3, 0.9, COL_ORANGE, "Retention", ["Holds", "Policies"], key_name="RET")
create_block(x3, 6.0, w3, 1.6, RGBColor(150, 0, 0), "RISE Services", ["Tokenization", "Masking"], key_name="RISE")

# Column 4
x4 = 8.3; w4 = 2.2
create_block(x4, 3.5, w4, 1.2, COL_GREEN, "Integration", ["Migration", "Bulk Export"], key_name="INT")
create_block(x4, 4.9, w4, 1.0, COL_RED, "Audit", ["Audit trails", "Compliance"], key_name="AUDIT")
create_block(x4, 6.1, w4, 1.5, COL_DARK_GREY, "Security", ["SSO", "RBAC"], key_name="SEC")

# Column 5
x5 = 10.7; w5 = 2.2
create_block(x5, 3.5, w5, 1.2, COL_YELLOW, "Repository Svcs", ["CRUD", "Routing"], key_name="REPO_SVC")
create_block(x5, 4.9, w5, 1.2, COL_ORANGE, "BPM Services", ["Workflow", "Orchestration"], key_name="BPM")
create_block(x5, 6.3, w5, 1.3, COL_YELLOW, "Relationship", ["Docs-Account", "Docs-Customer"], key_name="REL")

# --- 5. RIGHT SIDE (KAFKA) ---
create_block(13.7, 3.0, 1.8, 4.5, COL_GREY_BG, "Kafka Layer", ["Request events", "Doc events"], key_name="KAFKA")

# --- 6. BOTTOM ROW (REPOS) ---
y_bot = 8.2
repos = [
    ("FileNet", COL_BLUE_LIGHT), ("Documentum", COL_PURPLE), ("CM8", COL_TEAL),
    ("MARS", COL_ORANGE), ("Image Svcs", COL_BLUE_LIGHT), ("CMIS", COL_GREEN),
    ("CMOD", COL_BLUE_DARK), ("Legacy", COL_YELLOW)
]
for i, (name, col) in enumerate(repos):
    create_block(0.5 + (i * 1.3), y_bot, 1.1, 0.8, col, name, [], key_name=f"REPO_{i}")

# Add Caps Box
create_block(11.5, 8.2, 4.0, 0.8, COL_BLUE_DARK, "Add. Capabilities", ["Search", "Enrichment"], key_name="CAPS")


# --- CONNECTING THE LINES ---

# 1. Top Sources -> Firewall
for source in sources:
    add_connector(source[0], "PAA Firewalls")

# 2. Firewall -> Core
# Draw a connector from Firewall to the top of the Container
add_connector("PAA Firewalls", "CORE_CONTAINER")

# 3. Internal Flow (Left to Right approximation)
# Column 1 to Column 2
add_side_connector("API", "AI")
add_side_connector("UI", "TRANS")
add_side_connector("IVaaS", "CONV")

# Column 2 to Column 3
add_side_connector("AI", "DATA")
add_side_connector("TRANS", "RET")
add_side_connector("CONV", "RISE")

# Column 3 to Column 4
add_side_connector("DATA", "INT")
add_side_connector("RET", "AUDIT")
add_side_connector("RISE", "SEC")

# Column 4 to Column 5
add_side_connector("INT", "REPO_SVC")
add_side_connector("AUDIT", "BPM")
add_side_connector("SEC", "REL")

# 4. Output Flows
# BPM -> Kafka
add_side_connector("BPM", "KAFKA")

# Core -> Repositories
# We'll connect the bottom middle blocks to the repo line
for i in range(8):
    add_connector("SEC", f"REPO_{i}") # Logic: Security/Core feeds into Repos

# Core -> Capabilities
add_connector("REL", "CAPS")

# Save
prs.save('ICMP_Matrix_With_Lines.pptx')
print("Presentation saved with connectors.")
