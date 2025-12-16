from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create the presentation
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout

# Function to create a styled block
def create_block(left, top, width, height, color_rgb, title, content_list, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    
    # Fill Color
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color_rgb
    
    # Remove outline
    shape.line.fill.background()

    # Add Title
    text_frame = shape.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(255, 255, 255) # White text for headers
    p.alignment = PP_ALIGN.LEFT

    # Add Content
    for item in content_list:
        p = text_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.LEFT

# --- Colors (Approximated from image) ---
COL_BLUE_DARK = RGBColor(47, 85, 151)
COL_BLUE_LIGHT = RGBColor(68, 114, 196)
COL_PURPLE = RGBColor(112, 48, 160)
COL_TEAL = RGBColor(0, 176, 240)
COL_ORANGE = RGBColor(237, 125, 49)
COL_YELLOW = RGBColor(255, 192, 0)
COL_GREEN = RGBColor(146, 208, 80)
COL_RED = RGBColor(255, 100, 100)
COL_GREY_BG = RGBColor(242, 242, 242)
COL_DARK_GREY = RGBColor(89, 89, 89)

# --- HEADER ---
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(15), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "ICMP – Enterprise Content Management Capabilities Matrix"
p.font.bold = True
p.font.size = Pt(24)
p.font.color.rgb = COL_BLUE_DARK

p2 = tf.add_paragraph()
p2.text = "Document Ingestion, AI Intelligence, Processing, Retrieval, Repository Services, Security, BPM, and Notifications"
p2.font.size = Pt(14)

# --- TOP ROW (Ingestion Sources) ---
y_top = 1.5
w_small = 1.5
h_small = 0.6
gap = 0.2
start_x = 0.5

sources = [
    ("ICMP UI", COL_BLUE_LIGHT),
    ("Apigee API\nGateway", COL_PURPLE),
    ("NAS (Network\nAttached Storage)", COL_TEAL),
    ("Apache Kafka\nStreams", COL_ORANGE),
    ("Email\nIngestion", COL_YELLOW),
    ("Fax", COL_GREEN),
    ("Capture/KOFAX", COL_RED),
    ("Iron Mountain\nScanning", COL_RED)
]

for i, (name, col) in enumerate(sources):
    create_block(start_x + (i * (w_small + gap)), y_top, w_small, h_small, col, name, [])

# --- PAA Firewalls (Center) ---
create_block(6.5, 2.4, 3, 0.4, COL_DARK_GREY, "PAA Firewalls", [], MSO_SHAPE.RECTANGLE)

# --- MAIN CONTAINER (Core Processing) ---
# Create a large grey background rectangle
bg_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.0), Inches(13), Inches(4.5))
bg_shape.fill.solid()
bg_shape.fill.fore_color.rgb = COL_GREY_BG
bg_shape.line.color.rgb = COL_DARK_GREY

# Container Title
bg_tf = bg_shape.text_frame
bg_p = bg_tf.paragraphs[0]
bg_p.text = "ICMP Core Processing & Services Matrix"
bg_p.font.color.rgb = RGBColor(0,0,0)
bg_p.alignment = PP_ALIGN.CENTER
bg_p.font.bold = True

# --- INSIDE CONTAINER BLOCKS ---
# Column 1: API/UI
x_col1 = 0.7
w_col1 = 2.5
create_block(x_col1, 3.5, w_col1, 1.5, COL_BLUE_DARK, "API Layer", 
             ["Ingestion & validation", "Content retrieval via API", "Search, retrieve, export", "Metadata & property updates", "Supports Bulk Operations", "Log pattern masking"])

create_block(x_col1, 5.2, w_col1, 1.0, COL_BLUE_LIGHT, "ICMP UI Module", 
             ["Document upload", "Manual re-index & reprocess", "Bulk operations panel", "Operational dashboards"])

create_block(x_col1, 6.4, w_col1, 1.0, COL_BLUE_LIGHT, "IVaaS - Image Viewer", 
             ["View documents upstream", "Secure rendering", "Multiple formats (TIFF, PDF)", "Auth passthrough"])

# Column 2: Intelligence & Transformation
x_col2 = 3.4
w_col2 = 2.5
create_block(x_col2, 3.5, w_col2, 2.0, COL_BLUE_LIGHT, "ECM AI – iDocs Intelligence", 
             ["Document Classification", "Single/multi-context classification", "Document summarisation", "AI data extraction", "Signature detection"])

create_block(x_col2, 5.7, w_col2, 0.8, COL_PURPLE, "Transformation Services (TS)", 
             ["Metadata mapping", "Normalization rules", "Validation pipelines"])

create_block(x_col2, 6.7, w_col2, 0.7, COL_TEAL, "Conversion Services (CS)", 
             ["Format conversion (PDF, TIFF)", "Split & merge capabilities"])

# Column 3: Data & Retention
x_col3 = 6.1
w_col3 = 2.2
create_block(x_col3, 3.5, w_col3, 1.2, RGBColor(180, 50, 100), "Content Data Services", 
             ["OCR/ICR", "Data lift for managed content", "Push metadata to MongoDB"])

create_block(x_col3, 4.9, w_col3, 0.8, COL_ORANGE, "Retention Management (RM)", 
             ["Auto-apply retention", "Legal & audit holds"])

create_block(x_col3, 5.9, w_col3, 1.5, RGBColor(150, 0, 0), "ECM RISE Services", 
             ["PAN -> CRN tokenization", "Search by PAN or CRN", "Secure masking & retrieval"])

# Column 4: Integration & Security
x_col4 = 8.5
w_col4 = 2.2
create_block(x_col4, 3.5, w_col4, 1.2, COL_GREEN, "Integration Manager (IM)", 
             ["Large-scale ingestion", "Migration pipelines", "Retention application", "Bulk export"])

create_block(x_col4, 4.9, w_col4, 1.2, COL_RED, "Audit Module", 
             ["End-to-end audit trail", "Compliance lineage"])

create_block(x_col4, 6.3, w_col4, 1.1, COL_DARK_GREY, "Security Layer", 
             ["SSO-based authentication", "Dynamic security model", "Sensitive data masking"])

# Column 5: Repo & BPM
x_col5 = 10.9
w_col5 = 2.2
create_block(x_col5, 3.5, w_col5, 1.2, COL_YELLOW, "Repository Services (RS)", 
             ["CRUD operations", "Dynamic repository routing", "Repository abstraction layer"])

create_block(x_col5, 4.9, w_col5, 1.2, COL_ORANGE, "BPM Services", 
             ["Business data modifications", "Workflow orchestration", "Audit-compliant pipeline"])

create_block(x_col5, 6.3, w_col5, 1.1, COL_YELLOW, "Relationship Management", 
             ["Documents <-> Account", "Documents <-> Customer", "Documents <-> Deal", "Documents <-> Employee"])

# --- RIGHT SIDE (Kafka) ---
create_block(13.7, 3.0, 1.8, 4.5, COL_GREY_BG, "Kafka Notifications Layer", 
             ["Request-level events", "Document-level events", "Package-level events", "Configurable notification publishing"])

# --- BOTTOM ROW (Repositories) ---
y_bot = 8.0
repos = [
    ("FileNet", COL_BLUE_LIGHT),
    ("Documentum", COL_PURPLE),
    ("CM8", COL_TEAL),
    ("MARS", COL_ORANGE),
    ("Image\nServices", COL_BLUE_LIGHT),
    ("CMIS", COL_GREEN),
    ("CMOD", COL_BLUE_DARK),
    ("Legacy\nRepositories", COL_YELLOW)
]

w_repo = 1.2
for i, (name, col) in enumerate(repos):
    create_block(0.5 + (i * (w_repo + 0.1)), y_bot, w_repo, 0.8, col, name, [])

# Additional Capabilities (Bottom Right)
create_block(11.5, 8.0, 4.0, 0.8, COL_BLUE_DARK, "Additional Capabilities", 
             ["Search & retrieve", "Export", "Intelligent metadata enrichment", "High-scale migrations", "AI-driven document intelligence"])

# Save file
prs.save('ICMP_Matrix.pptx')
print("Presentation saved as ICMP_Matrix.pptx")
